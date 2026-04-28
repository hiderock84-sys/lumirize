from __future__ import annotations

import math
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.util import Cm


ROOT = Path(__file__).resolve().parent
ASSET_DIR = ROOT / "artifacts" / "official"
OUT_DIR = ROOT / "artifacts" / "output"
SLIDE_DIR = ROOT / "artifacts" / "designed_slides"
OUT_DIR.mkdir(parents=True, exist_ok=True)
SLIDE_DIR.mkdir(parents=True, exist_ok=True)

W, H = 1920, 1080
NAVY = (7, 24, 44)
INK = (13, 36, 61)
BLUE = (35, 108, 255)
CYAN = (0, 181, 203)
GREEN = (17, 184, 135)
ORANGE = (245, 158, 24)
RED = (238, 80, 87)
PURPLE = (119, 86, 255)
PINK = (237, 83, 145)
YELLOW = (255, 213, 74)
BG = (236, 243, 249)
WHITE = (255, 255, 255)
TEXT = (27, 42, 58)
MUTED = (82, 101, 121)
LINE = (199, 215, 229)
PALE = (247, 251, 253)

FONT_REG = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"
FONT_SERIF = "/usr/share/fonts/opentype/noto/NotoSerifCJK-Bold.ttc"


def font(size: int, bold: bool = False, serif: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_SERIF if serif else (FONT_BOLD if bold else FONT_REG), size)


def wrap_text(value: str, size: int, width: int) -> list[str]:
    chars = max(5, int(width / (size * 0.58)))
    lines: list[str] = []
    for paragraph in str(value).split("\n"):
        if not paragraph:
            lines.append("")
        else:
            lines.extend(textwrap.wrap(paragraph, chars, break_long_words=True, replace_whitespace=False))
    return lines


def draw_text(d: ImageDraw.ImageDraw, x: int, y: int, value: str, size: int, fill=TEXT,
              bold: bool = False, width: int | None = None, line_gap: int = 7,
              anchor: str | None = None, serif: bool = False):
    f = font(size, bold, serif)
    if width is None:
        d.text((x, y), value, font=f, fill=fill, anchor=anchor)
        return
    yy = y
    for line in wrap_text(value, size, width):
        d.text((x, yy), line, font=f, fill=fill)
        yy += size + line_gap


def rounded(img: Image.Image, box, radius=18, fill=WHITE, outline=LINE, width=2, shadow=True):
    x1, y1, x2, y2 = box
    if shadow:
        layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        ld = ImageDraw.Draw(layer)
        ld.rounded_rectangle((x1 + 8, y1 + 10, x2 + 8, y2 + 10), radius, fill=(21, 52, 84, 42))
        layer = layer.filter(ImageFilter.GaussianBlur(10))
        img.alpha_composite(layer)
    d = ImageDraw.Draw(img)
    d.rounded_rectangle(box, radius, fill=fill, outline=outline, width=width)


def fit_cover(path: Path, size: tuple[int, int]) -> Image.Image:
    im = Image.open(path).convert("RGB")
    iw, ih = im.size
    tw, th = size
    ratio = max(tw / iw, th / ih)
    im = im.resize((int(iw * ratio), int(ih * ratio)), Image.Resampling.LANCZOS)
    iw, ih = im.size
    return im.crop(((iw - tw) // 2, (ih - th) // 2, (iw + tw) // 2, (ih + th) // 2))


def paste_photo(img: Image.Image, path: Path, x: int, y: int, w: int, h: int,
                radius: int = 18, overlay=(0, 0, 0, 0), outline=WHITE):
    im = fit_cover(path, (w, h)).convert("RGBA")
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, w, h), radius, fill=255)
    img.paste(im, (x, y), mask)
    if overlay[3]:
        ov = Image.new("RGBA", (w, h), overlay)
        img.paste(ov, (x, y), mask)
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((x, y, x + w, y + h), radius, outline=outline, width=2)


def pill(d: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int, label: str,
         fill=BLUE, color=WHITE, size=18):
    d.rounded_rectangle((x, y, x + w, y + h), h // 2, fill=fill)
    d.text((x + w // 2, y + h // 2), label, font=font(size, True), fill=color, anchor="mm")


def base_slide(n: int, title: str, subtitle: str, section: str, accent=CYAN, total=24) -> Image.Image:
    img = Image.new("RGBA", (W, H), BG + (255,))
    d = ImageDraw.Draw(img)
    # Regular 12-column grid. Visible enough to make any remaining negative space intentional.
    margin, gutter = 54, 18
    col_w = (W - margin * 2 - gutter * 11) // 12
    for i in range(12):
        x = margin + i * (col_w + gutter)
        d.rectangle((x, 120, x + col_w, H - 62), fill=(255, 255, 255, 32))
    d.rectangle((0, 0, W, 120), fill=NAVY)
    d.rectangle((0, 120, W, 126), fill=accent)
    d.rounded_rectangle((54, 30, 106, 82), 14, fill=(18, 58, 92))
    d.text((80, 56), f"{n:02d}", font=font(18, True), fill=WHITE, anchor="mm")
    draw_text(d, 132, 24, title, 34, WHITE, True, width=920)
    draw_text(d, 134, 72, subtitle, 17, (194, 217, 234), width=1050)
    pill(d, 1510, 32, 310, 48, section, (238, 244, 249), NAVY, 18)
    d.rectangle((0, H - 48, W, H), fill=(246, 250, 253))
    draw_text(d, 54, H - 36, "SAGAMIHARA DARC | CHANGE YOUR LIFE!", 16, MUTED, True)
    d.text((1718, H - 28), f"{n}/{total}", font=font(17, True), fill=NAVY, anchor="ra")
    d.rounded_rectangle((1738, H - 29, 1872, H - 19), 5, fill=(209, 222, 232))
    d.rounded_rectangle((1738, H - 29, 1738 + int(134 * n / total), H - 19), 5, fill=accent)
    return img


def feature_card(img: Image.Image, x: int, y: int, w: int, h: int, title: str, body: list[str] | str,
                 accent=BLUE, title_size=24, body_size=18, dark=False):
    fill = INK if dark else WHITE
    text_color = WHITE if dark else TEXT
    muted = (218, 235, 246) if dark else MUTED
    rounded(img, (x, y, x + w, y + h), 18, fill, (56, 92, 122) if dark else LINE)
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((x, y, x + w, y + 10), 6, fill=accent)
    draw_text(d, x + 24, y + 26, title, title_size, WHITE if dark else NAVY, True, width=w - 48)
    lines = body if isinstance(body, list) else [body]
    yy = y + 76
    for item in lines:
        d.rounded_rectangle((x + 24, yy + 4, x + 40, yy + 20), 5, fill=accent)
        draw_text(d, x + 54, yy - 1, item, body_size, muted if dark else text_color, width=w - 78, line_gap=5)
        yy += max(36, body_size + 20)


def metric_card(img: Image.Image, x: int, y: int, w: int, h: int, label: str, value: str,
                note: str, accent=CYAN):
    rounded(img, (x, y, x + w, y + h), 18, INK, (48, 86, 118))
    d = ImageDraw.Draw(img)
    d.rectangle((x, y, x + w, y + 9), fill=accent)
    draw_text(d, x + 24, y + 24, label, 18, (188, 213, 231), True, width=w - 48)
    draw_text(d, x + 24, y + 66, value, 48, WHITE, True, width=w - 48)
    draw_text(d, x + 24, y + 132, note, 17, (204, 224, 237), width=w - 48)


def close_band(img: Image.Image, message: str, accent=CYAN):
    d = ImageDraw.Draw(img)
    rounded(img, (54, 910, 1866, 988), 18, NAVY, (42, 84, 116))
    d.rounded_rectangle((78, 932, 118, 970), 10, fill=accent)
    d.text((98, 951), "✓", font=font(23, True), fill=WHITE, anchor="mm")
    draw_text(d, 140, 930, message, 24, WHITE, True, width=1600)


def section_title(d: ImageDraw.ImageDraw, x: int, y: int, label: str, accent=CYAN):
    d.rounded_rectangle((x, y, x + 14, y + 50), 7, fill=accent)
    draw_text(d, x + 26, y + 8, label, 25, NAVY, True)


PHOTOS = {
    "main": "facility_main.jpg",
    "business1": "business01.jpg",
    "business2": "business02.jpg",
    "base1": "base01.jpg",
    "base2": "base02.jpg",
    "base3": "base03.jpg",
    "base4": "base04.jpg",
    "base5": "base05.jpg",
    "base6": "base06.jpg",
    "base7": "base07.jpg",
    "base8": "base08.jpg",
    "outline1": "outline01.jpg",
    "outline3": "outline03.jpg",
    "outline4": "outline04.jpg",
}


def p(key: str) -> Path:
    return ASSET_DIR / PHOTOS[key]


def slide_01() -> Image.Image:
    img = Image.new("RGBA", (W, H), NAVY + (255,))
    d = ImageDraw.Draw(img)
    paste_photo(img, p("main"), 1050, 0, 870, 1080, 0, (7, 24, 44, 92), outline=NAVY)
    d.polygon([(0, 0), (1210, 0), (860, 1080), (0, 1080)], fill=(7, 24, 44, 246))
    for y in range(80, 900, 72):
        d.line((0, y, 980, y + 240), fill=(0, 181, 203, 48), width=2)
    draw_text(d, 88, 66, "一般社団法人 相模原ダルク", 22, (202, 222, 236))
    pill(d, 1260, 48, 430, 54, "OFFICIAL PRESENTATION 2026", (18, 60, 96), WHITE, 20)
    draw_text(d, 98, 205, "CHANGE", 78, WHITE, True)
    draw_text(d, 98, 302, "YOUR LIFE!", 78, GREEN, True)
    draw_text(d, 650, 174, "DARC", 190, (230, 247, 244), True, serif=True)
    draw_text(d, 100, 455, "依存症からの回復支援施設\n新・施設紹介資料", 44, WHITE, True, width=760, line_gap=14)
    draw_text(d, 104, 585, "人としての尊厳を大切に、生き直す力を支える。\n安心できる共同の場・誠実な対話・専門機関連携で、回復の継続を支えます。", 27, (220, 237, 247), width=840, line_gap=10)
    feature_card(img, 96, 730, 840, 118, "資料の目的", ["本人・家族・関係機関に、相模原ダルクの支援全体像と相談導線を明確に伝える"], CYAN, 25, 20)
    for i, m in enumerate([("支援実績", "400+", "名以上", BLUE), ("卒業継続", "90%+", "回復継続", GREEN), ("見守り", "24h", "安心体制", ORANGE), ("受入", "70名", "程度", PURPLE)]):
        metric_card(img, 92 + i * 440, 890, 402, 152, *m)
    return img


AGENDA = [
    ("01", "社会課題", "依存症情勢・相談の入口", BLUE),
    ("02", "理念", "尊厳・共同の場・生き直す力", CYAN),
    ("03", "支援モデル", "相談から地域定着まで", GREEN),
    ("04", "事業", "障害福祉・入寮・相談・連携・啓発", ORANGE),
    ("05", "施設群", "公式写真で見る拠点ネットワーク", PURPLE),
    ("06", "プログラム", "デイ・ナイト・個別・5ステージ", BLUE),
    ("07", "疾患理解", "依存症・交差依存・PAWS", RED),
    ("08", "家族支援", "CRAFT・家族会・相談導線", GREEN),
    ("09", "連携", "医療・行政・司法・地域", CYAN),
    ("10", "行動", "問い合わせ・見学・入所検討", ORANGE),
]


def slide_02() -> Image.Image:
    img = base_slide(2, "全体構成｜10テーマで理解する相模原ダルク", "読み手が迷わないよう、相談から回復定着までを一つのストーリーに再構成", "CONTENTS", BLUE)
    d = ImageDraw.Draw(img)
    for i, (num, title, sub, col) in enumerate(AGENDA):
        x = 70 + (i % 2) * 895
        y = 160 + (i // 2) * 142
        rounded(img, (x, y, x + 825, y + 112), 18, WHITE, LINE)
        pill(d, x + 26, y + 24, 68, 64, num, col, WHITE, 23)
        draw_text(d, x + 122, y + 24, title, 29, NAVY, True)
        draw_text(d, x + 122, y + 68, sub, 20, MUTED, width=560)
        d.text((x + 760, y + 56), "→", font=font(32, True), fill=col, anchor="mm")
    close_band(img, "新資料では、ページごとに結論・根拠・次アクションを明確化し、余白も12カラムグリッド上で規則的に管理します。", BLUE)
    return img


def slide_03() -> Image.Image:
    img = base_slide(3, "社会課題｜依存症は本人だけで抱えられない病気", "アルコール・薬物・ギャンブルに共通する症状と相談サインを整理", "ISSUE", RED)
    d = ImageDraw.Draw(img)
    cols = [
        ("アルコール依存症", ["飲酒量や頻度をコントロールできない", "健康問題があっても飲酒を続ける", "仕事や家庭より飲酒を優先", "家族関係が悪化する"], BLUE),
        ("薬物依存症", ["薬物使用をやめようとしても繰り返す", "身体・精神への悪影響を軽視", "家庭や仕事より薬物使用を優先", "逮捕・訴訟など法的課題に発展"], GREEN),
        ("ギャンブル依存症", ["賭け金や時間をコントロールできない", "借金・多重債務につながる", "負けを取り戻そうと繰り返す", "窃盗・横領など刑事課題に発展"], ORANGE),
    ]
    for i, (title, lines, col) in enumerate(cols):
        feature_card(img, 70 + i * 600, 160, 550, 520, title, lines, col, 30, 22)
    rounded(img, (70, 720, 1840, 870), 20, NAVY, (42, 84, 116))
    draw_text(d, 110, 752, "公式サイト共通メッセージ", 30, WHITE, True)
    draw_text(d, 110, 804, "依存症は意思の弱さではなく、WHOも認める病気です。正しい理解と専門的な治療、そして一人で抱え込まない継続支援が必要です。", 26, (220, 238, 248), width=1640)
    close_band(img, "初期段階の違和感、家族関係の悪化、生活バランスの崩れが見えた時点で相談へつなぐことが重要です。", RED)
    return img


def slide_04() -> Image.Image:
    img = base_slide(4, "疾患理解｜依存症に共通する回復阻害サイクル", "コントロール困難、強い欲求、周囲への影響を一枚で把握", "DISEASE MODEL", ORANGE)
    d = ImageDraw.Draw(img)
    cx, cy = 960, 455
    d.ellipse((cx - 235, cy - 235, cx + 235, cy + 235), outline=(183, 204, 222), width=10)
    nodes = [
        ("一度始めると\n止めにくい", 960, 185, BLUE),
        ("健康・仕事・家庭に\n悪影響が出る", 1320, 455, RED),
        ("やめる決意をする", 960, 725, GREEN),
        ("強い欲求で\n継続が難しい", 600, 455, ORANGE),
    ]
    for title, x, y, col in nodes:
        rounded(img, (x - 185, y - 64, x + 185, y + 64), 22, WHITE, LINE)
        pill(d, x - 160, y - 22, 44, 44, "", col)
        draw_text(d, x - 100, y - 38, title, 24, NAVY, True, width=250)
    d.rounded_rectangle((790, 370, 1130, 540), 32, fill=NAVY)
    d.text((960, 455), "依存症\nサイクル", font=font(42, True), fill=WHITE, anchor="mm")
    feature_card(img, 70, 150, 380, 620, "3つの特徴", ["2種類の離脱症状：短期と長期の不調が再発の引き金になる", "優先順位の変化：依存対象が生活の中心になる", "周りを巻き込む：家族や周囲にも症状が出る"], PURPLE, 29, 21)
    feature_card(img, 1470, 150, 380, 620, "治療の前提", ["依存症そのものを治療対象として捉える", "医療機関と連携して離脱症状に対応", "仲間と支援につながり続ける環境をつくる"], GREEN, 29, 21)
    close_band(img, "問題行動だけを責めるのではなく、依存症サイクルを断ち切る支援環境を設計することが回復の出発点です。", ORANGE)
    return img


def slide_05() -> Image.Image:
    img = base_slide(5, "理念｜人としての尊厳を大切に、生き直す力を支える", "公式サイト掲載の理念・使命を、支援現場で使える判断軸として再構成", "PHILOSOPHY", CYAN)
    paste_photo(img, p("outline1"), 70, 160, 560, 590, 22, (7, 24, 44, 60))
    feature_card(img, 670, 160, 580, 260, "PHILOSOPHY｜理念", ["一人ひとりを尊重する関係性", "安心して自分を表現できる共同の場", "誠実な関わりと対話で生き直す力を育てる"], BLUE, 30, 22)
    feature_card(img, 1270, 160, 580, 260, "MISSION｜使命", ["依存症からの回復支援を通じて社会に貢献", "恐れや圧力ではなく安心できる環境で支える", "施設・地域・社会に共同の場を作る"], GREEN, 30, 22)
    rows = [("事業者", "一般社団法人 相模原ダルク"), ("代表理事", "田中 秀泰"), ("設立", "2014年3月3日"), ("所在地", "相模原市中央区千代田3-3-20")]
    for i, (k, v) in enumerate(rows):
        x, y = 670 + (i % 2) * 600, 470 + (i // 2) * 110
        rounded(img, (x, y, x + 580, y + 82), 16, WHITE, LINE)
        d = ImageDraw.Draw(img)
        draw_text(d, x + 28, y + 20, k, 20, CYAN, True)
        draw_text(d, x + 150, y + 20, v, 23, NAVY, True, width=390)
    close_band(img, "理念は飾りではなく、毎日の共同生活・対話・家族支援・地域連携のすべてに通底する運営基準です。", CYAN)
    return img


def slide_06() -> Image.Image:
    img = base_slide(6, "支援モデル｜相談から社会参加までを一気通貫でつなぐ", "本人・家族・関係機関が迷わないよう、支援の流れを見える化", "RECOVERY MODEL", GREEN)
    d = ImageDraw.Draw(img)
    steps = [
        ("01", "相談", "本人・家族・周囲から状況を丁寧に聴く", BLUE),
        ("02", "安全確保", "生活環境・緊急度・医療/司法課題を整理", CYAN),
        ("03", "入所/通所", "共同生活と日中活動で孤立をほどく", GREEN),
        ("04", "プログラム", "12ステップ・CBT・面談で再発予防", ORANGE),
        ("05", "社会参加", "就労・地域・家族関係へつなぎ直す", PURPLE),
    ]
    for i, (num, title, body, col) in enumerate(steps):
        x = 70 + i * 362
        rounded(img, (x, 210, x + 315, 620), 22, WHITE, LINE)
        pill(d, x + 28, 238, 72, 62, num, col, WHITE, 22)
        draw_text(d, x + 28, 330, title, 32, NAVY, True)
        draw_text(d, x + 28, 385, body, 22, TEXT, width=255, line_gap=9)
        if i < len(steps) - 1:
            d.line((x + 315, 410, x + 360, 410), fill=LINE, width=6)
            d.polygon([(x + 360, 410), (x + 342, 398), (x + 342, 422)], fill=LINE)
    feature_card(img, 70, 680, 860, 175, "本人への支援", ["生活リズム・仲間との関係・プログラム参加・段階的役割を通じて回復を継続"], BLUE, 27, 22)
    feature_card(img, 990, 680, 860, 175, "家族/関係機関への支援", ["相談、家族会、医療・福祉・司法との連携で本人だけに負担を集中させない"], GREEN, 27, 22)
    close_band(img, "回復支援は単発の面談ではなく、生活・関係・制度・地域をつなぐ継続的な設計です。", GREEN)
    return img


def slide_07() -> Image.Image:
    img = base_slide(7, "事業内容｜5つの事業を組み合わせた包括支援", "公式サイト掲載の事業を、相談者が理解しやすいサービスマップに整理", "BUSINESS", ORANGE)
    services = [
        ("障害福祉サービス", ["自立訓練（生活訓練）", "就労継続支援B型", "月額0円〜数千円程度"], BLUE),
        ("入寮事業", ["依存症治療専用ナイトケアハウス", "大型寮5ヶ所・個室寮9ヶ所", "月額176,000円（税込）"], GREEN),
        ("相談事業", ["本人・家族・周囲の方が相談可能", "秘密厳守", "初回無料／2回目以降3,000円"], CYAN),
        ("連携・協力事業", ["医療・行政・地域団体と連携", "回復支援ネットワーク形成", "監修・助言・運営協力"], PURPLE),
        ("予防・啓発事業", ["教育機関・行政・団体への講演", "地域イベントでの活動", "琉球太鼓エイサー演舞"], ORANGE),
    ]
    for i, (title, lines, col) in enumerate(services):
        x = 70 + (i % 3) * 600
        y = 160 + (i // 3) * 330
        feature_card(img, x, y, 550 if i < 3 else 850, 275, title, lines, col, 30, 22)
    close_band(img, "相談、生活訓練、入寮、就労、連携、啓発を別々にせず、本人の回復段階に応じて組み合わせます。", ORANGE)
    return img


def slide_08() -> Image.Image:
    img = base_slide(8, "入寮・住環境｜回復初期を支える安全な生活基盤", "集団療法の効果を高める大型寮と、自立準備の個室寮を組み合わせる", "RESIDENTIAL CARE", GREEN)
    paste_photo(img, p("business2"), 70, 160, 720, 540, 22, (7, 24, 44, 35))
    for i, m in enumerate([("大型寮", "5ヶ所", "集団療法の効果を最大化", BLUE), ("個室寮", "9ヶ所", "自立準備の練習環境", GREEN), ("費用", "176,000円", "月額・税込／生活保護は相談", ORANGE)]):
        metric_card(img, 835 + i * 335, 170, 300, 190, *m)
    feature_card(img, 835, 400, 1000, 300, "生活支援で重視すること", ["すべての寮が安全で安心して暮らせるよう配慮", "元の環境から距離を置き、強い欲求への対処を学ぶ", "共同生活を通じて孤立を解き、仲間と支え合う関係を作る"], CYAN, 31, 23)
    close_band(img, "回復初期は意思だけで乗り切る時期ではなく、安全な住環境と仲間の中で生活を整える時期です。", GREEN)
    return img


def slide_09() -> Image.Image:
    img = base_slide(9, "施設アトラス｜公式写真で見る相模原ダルクグループ", "施設写真を一覧化し、相談者が具体的な利用イメージを持てるページへ", "FACILITY ATLAS", PURPLE)
    photos = [
        ("デイケアセンター", "main"), ("OTC", "business1"), ("ナイトケア", "business2"), ("大和PCC", "base1"),
        ("町田RC", "base2"), ("愛川TC", "base3"), ("上溝HRC", "base4"), ("西門AC", "base5"),
        ("WPH", "base6"), ("関連施設", "base7"), ("関連施設", "base8"), ("外観/設備", "outline3"),
    ]
    d = ImageDraw.Draw(img)
    for i, (name, key) in enumerate(photos):
        x = 60 + (i % 4) * 455
        y = 150 + (i // 4) * 240
        paste_photo(img, p(key), x, y, 420, 170, 18, (0, 0, 0, 12), outline=WHITE)
        d.rounded_rectangle((x, y + 126, x + 420, y + 170), 0, fill=(7, 24, 44, 210))
        draw_text(d, x + 18, y + 134, name, 21, WHITE, True, width=380)
    close_band(img, "写真で具体像を示すことで、本人・家族が『どこで、どのように回復に取り組むのか』を想像しやすくします。", PURPLE)
    return img


def slide_10() -> Image.Image:
    img = base_slide(10, "プログラム構造｜デイ・ナイト・個別を統合する", "日中活動、夜間の安全、個別課題対応を組み合わせて生活全体を回復環境に変える", "PROGRAM", BLUE)
    paste_photo(img, p("business1"), 70, 160, 520, 250, 20, (7, 24, 44, 30))
    paste_photo(img, p("outline4"), 70, 440, 520, 250, 20, (7, 24, 44, 30))
    programs = [
        ("デイケア（通所）", ["ミーティング／プログラム", "12ステップ", "SAGARP（再発予防CBT）", "ワークショップ・アート", "感情・行動の振り返り"], BLUE),
        ("ナイトケア（寮）", ["生活習慣の改善", "金銭・食事の管理", "夜間の不安への対応", "24時間体制", "共同生活で孤立防止"], ORANGE),
        ("個別サポート", ["個別面談", "債務・司法相談", "家族支援", "就労準備", "医療・福祉への接続"], GREEN),
    ]
    for i, (title, lines, col) in enumerate(programs):
        feature_card(img, 635 + i * 410, 160, 380, 530, title, lines, col, 28, 20)
    close_band(img, "回復を日中だけでなく夜間・個別課題まで支えることで、再発リスクを生活全体から下げます。", BLUE)
    return img


def slide_11() -> Image.Image:
    img = base_slide(11, "5ステージ制｜回復を可視化するロードマップ", "段階的な目標設定と役割付与で、自立と社会参加へ進む", "5 STAGES", PURPLE)
    stages = [
        ("MEMBER", "0-3ヶ月", ["基礎知識", "共同生活", "生活習慣", "予定を守る"], BLUE),
        ("SUPPORT", "3-6ヶ月", ["役割付与", "振り返り", "感情の言語化", "再発リスク把握"], ORANGE),
        ("TRAINEE", "6-12ヶ月", ["リーダー練習", "他者支援", "外部活動", "生活管理"], GREEN),
        ("CHIEF", "12-18ヶ月", ["後輩支援", "相談補助", "就労/通院", "地域生活"], PURPLE),
        ("MANAGER", "18ヶ月-", ["運営補助", "社会参加", "再発予防計画", "卒業後計画"], PINK),
    ]
    d = ImageDraw.Draw(img)
    for i, (name, period, lines, col) in enumerate(stages):
        x = 55 + i * 374
        rounded(img, (x, 170, x + 340, 690), 20, WHITE, LINE)
        d.rounded_rectangle((x + 22, 194, x + 318, 260), 18, fill=col)
        d.text((x + 170, 218), name, font=font(25, True), fill=WHITE, anchor="mm")
        d.text((x + 170, 244), period, font=font(18, True), fill=(235, 246, 250), anchor="mm")
        yy = 300
        for line in lines:
            pill(d, x + 30, yy, 32, 32, "", col, WHITE, 10)
            draw_text(d, x + 76, yy - 2, line, 22, TEXT, True, width=230)
            yy += 78
    close_band(img, "ステージ制は、本人が『今どこにいて、次に何を目指すか』を把握し、成功体験を積み上げるための仕組みです。", PURPLE)
    return img


def slide_12() -> Image.Image:
    img = base_slide(12, "依存症の理解｜7つの特徴を支援設計に落とす", "病気として理解することで、責める支援から、つながる支援へ変える", "ADDICTION", RED)
    features = [
        ("一次性", "原因は依存症そのもの", BLUE), ("慢性", "継続的なケアが必要", GREEN), ("進行性", "放置すると悪化", ORANGE), ("死亡率", "事故・自殺等のリスク", RED),
        ("性格変化", "関係性が崩れる", PURPLE), ("対象移行", "交差依存が起きる", PINK), ("巻き込み", "家族・周囲にも影響", CYAN),
    ]
    for i, (title, body, col) in enumerate(features):
        x = 70 + (i % 4) * 445
        y = 165 + (i // 4) * 250
        feature_card(img, x, y, 410, 205, title, [body, "理解が支援の出発点"], col, 27, 19)
    feature_card(img, 70, 720, 1770, 150, "公式サイトの共通メッセージ", ["依存症は『正しい理解』と『専門的な治療』が必要な特殊な病気です。本人だけでやめ続けることは難しく、他者の力を借りることが大切です。"], NAVY, 28, 24, dark=True)
    close_band(img, "病気として理解すると、本人を責めるのではなく、再発リスクを下げる環境づくりに焦点が移ります。", RED)
    return img


def slide_13() -> Image.Image:
    img = base_slide(13, "交差依存｜対象が変わっても依存症は続く", "薬物・アルコール・ギャンブル・処方薬・ネット等を横断して捉える", "CROSS ADDICTION", ORANGE)
    d = ImageDraw.Draw(img)
    cx, cy = 960, 455
    d.ellipse((cx - 240, cy - 240, cx + 240, cy + 240), outline=(174, 197, 216), width=10)
    nodes = [("本命を止める", 960, 185, BLUE), ("苦しみが増える", 1330, 455, ORANGE), ("別対象へ移る", 960, 725, GREEN), ("本命へ戻る", 590, 455, PURPLE)]
    for title, x, y, col in nodes:
        rounded(img, (x - 175, y - 58, x + 175, y + 58), 20, WHITE, LINE)
        pill(d, x - 150, y - 20, 40, 40, "", col)
        draw_text(d, x - 95, y - 24, title, 24, NAVY, True, width=230)
    d.rounded_rectangle((800, 390, 1120, 520), 30, fill=NAVY)
    d.text((960, 455), "依存対象\n移行サイクル", font=font(32, True), fill=WHITE, anchor="mm")
    feature_card(img, 70, 170, 390, 560, "依存のタイプ", ["物質依存：酒・薬物", "プロセス依存：ギャンブル・ネット", "関係依存：共依存など", "対象が変わっても根本課題は残る"], RED, 28, 21)
    feature_card(img, 1460, 170, 390, 560, "支援の原則", ["対象だけでなく生活全体を見る", "代替行動を設計する", "家族・支援者と共有する", "早期に専門支援へつなぐ"], GREEN, 28, 21)
    close_band(img, "交差依存では『何をやめるか』だけではなく、『何につながり、どう生活を整えるか』が重要です。", ORANGE)
    return img


def slide_14() -> Image.Image:
    img = base_slide(14, "PAWS｜長期離脱症状を見越した支援", "回復初期から中長期にかけて起こる波を、チームで乗り越える", "PAWS", BLUE)
    symptoms = [("睡眠", "不眠・過眠・中途覚醒", BLUE), ("ストレス", "音・視線・人混みに過敏", ORANGE), ("記憶", "覚えられない・忘れる", GREEN), ("感情", "怒り・落ち込み・不安定", RED), ("身体", "疲労・めまい・不調", PURPLE), ("思考", "集中力低下・判断ミス", PINK)]
    for i, (title, body, col) in enumerate(symptoms):
        x = 70 + (i % 3) * 600
        y = 165 + (i // 3) * 245
        feature_card(img, x, y, 550, 200, title, [body, "予定・睡眠・食事・相談で波を整える"], col, 29, 20)
    d = ImageDraw.Draw(img)
    rounded(img, (70, 695, 1840, 840), 20, WHITE, LINE)
    draw_text(d, 110, 725, "回復プロセスの考え方", 30, NAVY, True)
    d.rounded_rectangle((500, 775, 1640, 804), 14, fill=(207, 221, 232))
    d.rounded_rectangle((500, 775, 840, 804), 14, fill=BLUE)
    d.rounded_rectangle((840, 775, 1210, 804), 14, fill=ORANGE)
    d.rounded_rectangle((1210, 775, 1640, 804), 14, fill=GREEN)
    draw_text(d, 500, 817, "急性期", 20, BLUE, True)
    draw_text(d, 840, 817, "波が出る時期", 20, ORANGE, True)
    draw_text(d, 1210, 817, "安定化", 20, GREEN, True)
    close_band(img, "PAWSを知らないと再発を本人の失敗として捉えがちです。症状の波を前提にした支援設計が必要です。", BLUE)
    return img


def slide_15() -> Image.Image:
    img = base_slide(15, "家族支援｜CRAFTと家族会で関わり方を再設計する", "家族は回復の重要パートナーであり、同時に支援されるべき当事者です", "FAMILY SUPPORT", GREEN)
    d = ImageDraw.Draw(img)
    steps = [("1", "理解・学ぶ"), ("2", "相談・家族会"), ("3", "関わりを変える"), ("4", "境界線を引く"), ("5", "継続する")]
    for i, (num, title) in enumerate(steps):
        x = 70 + i * 350
        rounded(img, (x, 170, x + 305, 300), 20, WHITE, LINE)
        pill(d, x + 28, 204, 58, 58, num, GREEN, WHITE, 22)
        draw_text(d, x + 104, 204, title, 25, NAVY, True, width=180)
    feature_card(img, 70, 350, 850, 320, "家族支援で扱うテーマ", ["依存症の正しい理解", "本人との距離の取り方", "境界線と安全確保", "相談・入所・治療へのつなぎ方", "家族自身の回復"], BLUE, 29, 22)
    feature_card(img, 980, 350, 850, 320, "援助の質を高めるポイント", ["原因探しをしすぎない", "責めない・さばかない", "回復行動を強化する", "タイミングを見て伝える", "家族だけで抱え込まない"], ORANGE, 29, 22)
    close_band(img, "家族支援は、本人を説得する技術ではなく、家族自身の安全・理解・関係性を整える支援です。", GREEN)
    return img


def slide_16() -> Image.Image:
    img = base_slide(16, "相談支援｜初期段階から入所までを丸ごと支える", "公式サイト掲載のサポートを、相談者にわかりやすい導線へ整理", "CONSULTATION", CYAN)
    supports = [
        ("初期段階なんでも相談", ["生活に支障が出てきた", "本人・家族・周囲の方が相談可能", "早期相談で回復率を高める"], BLUE),
        ("入所まるごとサポート", ["入所を考えているが動けない", "抱えている課題を整理", "治療開始まで支援"], GREEN),
        ("刑事裁判支援", ["逮捕・起訴・裁判への対応", "弁護士と協力", "治療を望む本人を支援"], ORANGE),
        ("債務整理支援", ["ギャンブル依存で多重債務", "弁護士・司法書士と連携", "治療と借金問題を同時支援"], PURPLE),
    ]
    for i, (title, lines, col) in enumerate(supports):
        x = 70 + (i % 2) * 900
        y = 165 + (i // 2) * 300
        feature_card(img, x, y, 850, 245, title, lines, col, 30, 22)
    close_band(img, "相談内容は秘密厳守。本人・家族・友人知人・関係機関など、どなたからの相談でも入口を作ります。", CYAN)
    return img


def slide_17() -> Image.Image:
    img = base_slide(17, "連携ネットワーク｜施設内で完結させない支援", "医療・行政・司法・地域とつながることで、生活再建の選択肢を広げる", "NETWORK", PURPLE)
    nets = [
        ("医療", ["解毒入院", "薬物療法", "精神科・身体面治療", "退院後の受け皿"], BLUE),
        ("行政・福祉", ["障害福祉サービス", "生活保護相談", "地域生活支援", "制度利用調整"], GREEN),
        ("司法・専門職", ["刑事裁判支援", "保護観察", "債務整理", "弁護士・司法書士連携"], ORANGE),
        ("地域・教育", ["予防啓発", "講演活動", "地域イベント", "学校・団体連携"], PURPLE),
    ]
    for i, (title, lines, col) in enumerate(nets):
        x = 70 + (i % 2) * 900
        y = 165 + (i // 2) * 300
        feature_card(img, x, y, 850, 245, title, lines, col, 30, 22)
    close_band(img, "回復支援は施設内で閉じません。医療・福祉・司法・地域の接続が、再発予防と社会参加を支えます。", PURPLE)
    return img


def slide_18() -> Image.Image:
    img = base_slide(18, "予防・啓発｜地域に回復の理解を広げる", "講演・研修・活動レポート・地域イベントを通じて、偏見を減らし相談につなぐ", "PREVENTION", ORANGE)
    paste_photo(img, p("business1"), 70, 160, 520, 300, 22, (7, 24, 44, 28))
    paste_photo(img, p("outline3"), 70, 500, 520, 300, 22, (7, 24, 44, 28))
    feature_card(img, 635, 160, 570, 300, "予防・啓発活動", ["地域の教育機関・行政機関・団体への講演", "依存症の予防・啓発に関する情報発信", "地域イベントでの琉球太鼓エイサー演舞"], BLUE, 30, 22)
    feature_card(img, 1260, 160, 570, 300, "相談につながる効果", ["依存症を隠すものから相談できるものへ", "家族や関係者が早期に気づく", "地域の理解が本人の孤立を減らす"], GREEN, 30, 22)
    feature_card(img, 635, 500, 1195, 300, "社会的メッセージ", ["依存症は回復できる病気です。正しい理解と専門的な支援が届くほど、本人と家族の選択肢は増えます。地域の中に回復を支える共同の場を広げていきます。"], ORANGE, 30, 24)
    close_band(img, "予防啓発は、依存症の発生を減らすだけでなく、相談につながる社会的な入口を増やす活動です。", ORANGE)
    return img


def slide_19() -> Image.Image:
    img = base_slide(19, "実績・運営基盤｜数字で見る支援体制", "支援実績、継続率、受入規模、相談体制を一枚で把握", "PERFORMANCE", BLUE)
    for i, m in enumerate([("支援実績", "400+", "名以上の回復支援", BLUE), ("卒業継続", "90%+", "地域での継続を支える", GREEN), ("再使用率", "5%以下", "入所中のリスク低減", RED), ("受入体制", "70名", "大型寮・個室寮", PURPLE)]):
        metric_card(img, 70 + i * 455, 170, 410, 250, *m)
    rounded(img, (70, 480, 1840, 800), 22, WHITE, LINE)
    d = ImageDraw.Draw(img)
    draw_text(d, 110, 515, "支援成果を支える5要素", 32, NAVY, True)
    labels = [("相談入口", BLUE), ("安全な住環境", GREEN), ("仲間との対話", ORANGE), ("専門連携", PURPLE), ("卒業後支援", CYAN)]
    for i, (label, col) in enumerate(labels):
        x = 120 + i * 330
        pill(d, x, 610, 235, 58, label, col, WHITE, 23)
        if i < 4:
            d.line((x + 245, 639, x + 305, 639), fill=LINE, width=6)
            d.polygon([(x + 305, 639), (x + 286, 625), (x + 286, 653)], fill=LINE)
    close_band(img, "数字は単なる実績ではなく、相談・住環境・プログラム・連携を継続運用できる基盤を示しています。", BLUE)
    return img


def slide_20() -> Image.Image:
    img = base_slide(20, "利用イメージ｜相談から支援開始まで", "問い合わせ後の流れを具体化し、不安を減らす", "PROCESS", GREEN)
    steps = [
        ("1", "問い合わせ", "電話・フォームから相談"),
        ("2", "状況整理", "本人/家族/関係機関から聴き取り"),
        ("3", "見学・面談", "支援方針と安全性を確認"),
        ("4", "利用開始", "相談・通所・入寮・連携へ"),
        ("5", "継続支援", "プログラム・家族支援・地域定着"),
    ]
    d = ImageDraw.Draw(img)
    for i, (num, title, body) in enumerate(steps):
        x = 70 + i * 360
        rounded(img, (x, 205, x + 310, 440), 22, WHITE, LINE)
        pill(d, x + 28, 235, 58, 58, num, GREEN, WHITE, 22)
        draw_text(d, x + 28, 320, title, 29, NAVY, True, width=240)
        draw_text(d, x + 28, 370, body, 21, TEXT, width=250)
    feature_card(img, 70, 520, 850, 270, "相談できる人", ["本人", "家族", "友人・知人", "医療・福祉・司法などの関係機関"], BLUE, 30, 23)
    feature_card(img, 980, 520, 850, 270, "相談内容", ["入所したい", "見学がしたい", "相談がしたい", "講演・連携など各種依頼"], ORANGE, 30, 23)
    close_band(img, "最初から答えが決まっていなくても構いません。まず状況を整理し、次の一歩を一緒に決めます。", GREEN)
    return img


def slide_21() -> Image.Image:
    img = base_slide(21, "比較で理解する｜一人で抱える状態から支援につながる状態へ", "家族や関係者に伝わるよう、支援前後の変化を可視化", "BEFORE / AFTER", CYAN)
    before = ["問題を隠す・否認する", "家族が責める/抱え込む", "生活リズムが崩れる", "医療・司法・福祉が分断", "再発時に孤立する"]
    after = ["病気として理解する", "相談先と家族支援がある", "共同生活で生活を整える", "専門機関と連携する", "再発リスクを共有し対応する"]
    feature_card(img, 90, 170, 760, 600, "支援につながる前", before, RED, 34, 25)
    feature_card(img, 1070, 170, 760, 600, "支援につながった後", after, GREEN, 34, 25)
    d = ImageDraw.Draw(img)
    d.rounded_rectangle((880, 410, 1040, 520), 30, fill=NAVY)
    d.text((960, 465), "CHANGE", font=font(26, True), fill=WHITE, anchor="mm")
    d.polygon([(1040, 465), (1088, 438), (1088, 492)], fill=NAVY)
    close_band(img, "資料の役割は、本人・家族・関係者が『何を変えるべきか』を同じ地図で見られるようにすることです。", CYAN)
    return img


def slide_22() -> Image.Image:
    img = base_slide(22, "施設概要｜相談先として必要な基本情報", "連絡・所在地・営業時間・費用を1枚に集約", "OUTLINE", PURPLE)
    items = [
        ("事業者", "一般社団法人 相模原ダルク"),
        ("代表理事", "田中 秀泰"),
        ("設立", "2014年3月3日"),
        ("所在地", "〒252-0237 神奈川県相模原市中央区千代田3-3-20"),
        ("TEL/FAX", "TEL 042-707-0391 / FAX 042-707-0392"),
        ("営業時間", "平日 9:00-17:00 / 土・祝 9:00-14:00（日曜定休）"),
        ("相談費用", "初回無料 / 2回目以降 3,000円（税込）"),
        ("入寮費用", "月額 176,000円（税込）※生活保護の方は相談可"),
    ]
    d = ImageDraw.Draw(img)
    for i, (k, v) in enumerate(items):
        x = 70 + (i % 2) * 900
        y = 160 + (i // 2) * 160
        rounded(img, (x, y, x + 850, y + 115), 18, WHITE, LINE)
        draw_text(d, x + 28, y + 24, k, 22, PURPLE, True)
        draw_text(d, x + 190, y + 24, v, 22, NAVY, True, width=610)
    close_band(img, "事前予約によりデイケア横の駐車スペースが利用可能です。まずは電話またはフォームからご相談ください。", PURPLE)
    return img


def slide_23() -> Image.Image:
    img = base_slide(23, "次アクション｜説明後に確認する5つのこと", "聞いて終わりではなく、相談・見学・連携へつなげるためのチェックリスト", "ACTION", ORANGE)
    actions = [
        ("本人の状況", "依存対象、使用/飲酒/ギャンブル頻度、緊急性"),
        ("家族の状況", "同居可否、疲弊度、連絡可能な支援者"),
        ("医療・司法課題", "通院、逮捕・裁判、債務、生活保護等"),
        ("利用希望", "相談、見学、通所、入寮、家族会"),
        ("連絡方法", "電話、メール、希望曜日・時間帯"),
    ]
    for i, (title, body) in enumerate(actions):
        x = 70 + (i % 3) * 600
        y = 170 + (i // 3) * 290
        feature_card(img, x, y, 550, 240, f"{i+1}. {title}", [body], [BLUE, GREEN, ORANGE, PURPLE, CYAN][i], 30, 23)
    close_band(img, "相談の目的は、完璧な答えを出すことではなく、孤立を止め、次にできる行動を一つ決めることです。", ORANGE)
    return img


def slide_24() -> Image.Image:
    img = base_slide(24, "お問い合わせ｜CHANGE YOUR LIFE!", "変われる。変われた仲間がいる。共に変わろう、共に進もう。", "CONTACT", BLUE)
    d = ImageDraw.Draw(img)
    draw_text(d, 440, 160, "CHANGE YOUR LIFE!", 66, NAVY, True)
    draw_text(d, 485, 250, "相模原ダルクは、依存症からの回復を全力で支えます。", 30, MUTED, True, width=980)
    feature_card(img, 80, 360, 520, 250, "電話でのご相談", ["042-707-0391", "平日 9:00-17:00", "土・祝 9:00-14:00"], BLUE, 31, 24)
    feature_card(img, 650, 360, 520, 250, "お問い合わせフォーム", ["https://s-darc.com", "24時間受付可能", "本人・家族・関係者から相談可"], GREEN, 31, 24)
    feature_card(img, 1220, 360, 520, 250, "所在地", ["〒252-0237", "相模原市中央区千代田3-3-20", "日曜定休"], ORANGE, 31, 24)
    rounded(img, (360, 690, 1560, 860), 28, BLUE, (67, 128, 255))
    draw_text(d, 430, 725, "お電話でのご相談", 30, WHITE, True)
    draw_text(d, 430, 775, "042-707-0391", 56, WHITE, True)
    draw_text(d, 430, 842, "初回相談無料・秘密厳守", 26, (230, 244, 255), True)
    close_band(img, "まず一歩、状況を聞かせてください。本人・ご家族・関係者、どなたからの相談でも構いません。", BLUE)
    return img


SLIDE_BUILDERS = [
    slide_01, slide_02, slide_03, slide_04, slide_05, slide_06, slide_07, slide_08,
    slide_09, slide_10, slide_11, slide_12, slide_13, slide_14, slide_15, slide_16,
    slide_17, slide_18, slide_19, slide_20, slide_21, slide_22, slide_23, slide_24,
]


def build_images() -> list[Path]:
    paths: list[Path] = []
    for idx, builder in enumerate(SLIDE_BUILDERS, 1):
        im = builder().convert("RGB")
        out = SLIDE_DIR / f"visual-slide-{idx:02d}.png"
        im.save(out, quality=96)
        paths.append(out)
    return paths


def build_pptx(paths: list[Path]) -> Path:
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    blank = prs.slide_layouts[6]
    for path in paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    out = OUT_DIR / "sagamihara_darc_professional_presentation_2026.pptx"
    root = ROOT / "Sagamihara_DARC_Professional_2026.pptx"
    prs.save(out)
    prs.save(root)
    return root


if __name__ == "__main__":
    pptx = build_pptx(build_images())
    print(pptx)
